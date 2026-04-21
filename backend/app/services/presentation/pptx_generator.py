import os
import copy
import random
import asyncio
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    from app.services.presentation.ai_generator import SlideData
except ImportError:
    from dataclasses import dataclass, field

    @dataclass
    class SlideData:
        index: int
        title: str
        bullets: list
        speaker_notes: str = ""
        slide_type: str = "content"
        image_keyword: str = ""
        raw_data: dict = field(default_factory=dict)

TEMPLATES_DIR = Path(__file__).parent / "templates"


def _get_template_path(template_index: Optional[int] = None) -> tuple[str, int]:
    if template_index is None:
        template_index = random.randint(1, 10)
    template_index = max(1, min(10, template_index))
    path = TEMPLATES_DIR / f"dizayn_ {template_index}.pptx"
    if not path.exists():
        for i in range(1, 11):
            fallback = TEMPLATES_DIR / f"dizayn_ {i}.pptx"
            if fallback.exists():
                return str(fallback), i
        raise FileNotFoundError(f"Templates papkasida shablon topilmadi: {TEMPLATES_DIR}")
    return str(path), template_index


def _clone_slide(prs: Presentation, template_slide) -> object:
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Shablon fondasidagi eski shakllarni tozalaymiz
    sp_tree = new_slide.shapes._spTree
    tags_to_remove = ['}sp', '}pic', '}grpSp', '}cxnSp']
    for sp in list(sp_tree):
        if any(sp.tag.endswith(t) for t in tags_to_remove):
            sp_tree.remove(sp)

    tmpl_sp_tree = template_slide.shapes._spTree
    for sp in tmpl_sp_tree:
        if any(sp.tag.endswith(t) for t in tags_to_remove):
            new_sp = copy.deepcopy(sp)
            sp_tree.append(new_sp)

    try:
        P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        tmpl_cSld = template_slide._element.find(f'{{{P_NS}}}cSld')
        new_cSld = new_slide._element.find(f'{{{P_NS}}}cSld')
        if tmpl_cSld is not None and new_cSld is not None:
            tmpl_bg = tmpl_cSld.find(f'{{{P_NS}}}bg')
            if tmpl_bg is not None:
                old_bg = new_cSld.find(f'{{{P_NS}}}bg')
                new_bg = copy.deepcopy(tmpl_bg)
                if old_bg is not None:
                    new_cSld.replace(old_bg, new_bg)
                else:
                    new_cSld.insert(0, new_bg)
    except Exception:
        pass

    try:
        for rel in template_slide.part.rels.values():
            if "image" in rel.reltype:
                new_slide.part.rels.get_or_add(rel.reltype, rel._target)
    except Exception:
        pass

    return new_slide


def _inject_title(slide, title: str):
    # Faqat 0-indeksdagi placeholderga yozamiz (Sarlavha)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title
            # Majburiy tarzda sarlavhani eng tepaga olib chiqish:
            try:
                shape.top = Inches(0.4)
                shape.height = Inches(0.8)
            except Exception:
                pass
            return
            
    # Agar 0-indeks topilmasa (shablon buzilgan bo'lsa), text framelarni qidiramiz
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    if text_shapes:
        text_shapes[0].text = title


def _inject_notes(slide, notes: str):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def get_theme_color(tmpl_idx: int) -> RGBColor:
    # Shablonga mos oq va qora ranglar:
    # 1: oq, 2: qora, 3: oq, 4: oq, 5: oq, 6: qora, 7: qora, 8: oq, 9: oq, 10: qora
    if tmpl_idx in [2, 6, 7, 10]:
        return RGBColor(0, 0, 0) # Qora
    else:
        return RGBColor(255, 255, 255) # Oq


def _render_slide_content(slide, sd: SlideData, img_path: Optional[str], slide_w, slide_h, tmpl_idx: int):
    # Birinchi navbatda sarlavha (idx ba'zan 0 bo'ladi) dan boshqa BARCHA body/text placeholderlarni o'chiramiz.
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx != 0 or sd.slide_type == "quote":
            sp = shape._element
            sp.getparent().remove(sp)
            
    # Asosiy koordinatalar (Sarlavha tepaga surilgani sababli margin_top 1.5 ga tushirildi)
    margin_x = Inches(0.5)
    margin_top = Inches(1.5) 
    margin_bottom = Inches(0.5)
    
    content_w = slide_w - margin_x * 2
    content_h = slide_h - margin_top - margin_bottom
    
    stype = sd.slide_type
    txt_color = get_theme_color(tmpl_idx)
    
    def add_text_box(x, y, w, h, text_lines, is_full_text=False):
        if not text_lines: return
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        if is_full_text:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        else:
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
        for i, line in enumerate(text_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(line)
            p.level = 0
            
            if is_full_text:
                p.font.size = Pt(28)
                p.alignment = PP_ALIGN.JUSTIFY
            else:
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                
            p.font.color.rgb = txt_color

    def add_image(x, y, w, h, img_p):
        if img_p and os.path.exists(img_p):
            try:
                # To'g'ri usul: add_picture orqali rasm qo'shish (python-pptx standart API)
                pic = slide.shapes.add_picture(img_p, x, y, w, h)
                print(f"[PptxGen] Rasm qo'shildi: {img_p}")
            except Exception as e:
                print(f"[PptxGen] Rasm xatosi: {e}")

    # ===== LAYOUTLAR =====
    if stype == "content_image_right":
        add_text_box(margin_x, margin_top, content_w * 0.45, content_h, sd.bullets)
        if img_path:
            add_image(margin_x + content_w * 0.5, margin_top, content_w * 0.5, content_h, img_path)
            
    elif stype == "content_image_left":
        if img_path:
            add_image(margin_x, margin_top, content_w * 0.5, content_h, img_path)
        add_text_box(margin_x + content_w * 0.55, margin_top, content_w * 0.45, content_h, sd.bullets)
        
    elif stype == "table":
        raw = sd.raw_data or {}
        table_dict = raw.get("table", {})
        headers = table_dict.get("headers", [])
        rows = table_dict.get("rows", [])
        
        total_rows = (1 if headers else 0) + len(rows)
        total_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
        
        if total_rows > 0 and total_cols > 0:
            shape = slide.shapes.add_table(total_rows, total_cols, margin_x, margin_top, content_w, content_h)
            tbl = shape.table
            row_idx = 0
            if headers:
                for col_idx, h_text in enumerate(headers):
                    if col_idx < len(tbl.columns):
                        tbl.cell(row_idx, col_idx).text = str(h_text)
                row_idx += 1
            for r in rows:
                if row_idx < len(tbl.rows):
                    for col_idx, c_text in enumerate(r):
                        if col_idx < len(tbl.columns):
                            tbl.cell(row_idx, col_idx).text = str(c_text)
                row_idx += 1
        else:
            add_text_box(margin_x, margin_top, content_w, content_h, sd.bullets, is_full_text=True)

    elif stype in ["chart_bar", "chart_pie", "chart_line"]:
        raw = sd.raw_data or {}
        chart_dict = raw.get("chart", {})
        data_points = chart_dict.get("data", [])
        insight = raw.get("insight", "")
        
        if data_points:
            chart_data = CategoryChartData()
            categories = []
            values = []
            for dp in data_points:
                categories.append(str(dp.get("label", "")))
                values.append(float(dp.get("value", 0)))
                
            chart_data.categories = categories
            chart_data.add_series("Ma'lumotlar", tuple(values))
            
            ctype = XL_CHART_TYPE.COLUMN_CLUSTERED
            if stype == "chart_pie": ctype = XL_CHART_TYPE.PIE
            elif stype == "chart_line": ctype = XL_CHART_TYPE.LINE
            
            # Diagrammani chizish (80% balandlikda)
            chart_h = content_h * 0.8 if insight else content_h
            try:
                slide.shapes.add_chart(ctype, margin_x, margin_top, content_w, chart_h, chart_data)
                
                # Insight xulosasini pastga yozamiz
                if insight:
                    add_text_box(margin_x, margin_top + chart_h + Inches(0.1), content_w, content_h * 0.15, [f"Xulosa: {insight}"])
            except Exception as e:
                print(f"[PptxGen] Chart xatosi: {e}")
                add_text_box(margin_x, margin_top, content_w, content_h, sd.bullets, is_full_text=True)
        else:
            add_text_box(margin_x, margin_top, content_w, content_h, sd.bullets, is_full_text=True)

    elif stype == "quote":
        raw = sd.raw_data or {}
        q_text = raw.get("quote", "")
        author = raw.get("author", "")
        
        text_lines = []
        if q_text: text_lines.append(f'"{q_text}"')
        if author: text_lines.append(f"— {author}")
        if not text_lines: text_lines = sd.bullets
        
        txBox = slide.shapes.add_textbox(margin_x, margin_top + Inches(0.5), content_w, Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36 if i == 0 else 24)
            p.font.italic = (i == 0)
            p.font.color.rgb = txt_color

    else:
        # Odatiy matn slayd (content, conclusion, section, agenda va title)
        if stype != "title":
            add_text_box(margin_x, margin_top, content_w, content_h, sd.bullets, is_full_text=True)


def _delete_slide(prs: Presentation, slide_index: int):
    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]


def _build_presentation(
    slides: list,
    output_path: str,
    template_index: Optional[int] = None,
    user_images: Optional[list] = None,
) -> str:
    template_path, actual_template_index = _get_template_path(template_index)
    print(f"[PptxGen] Shablon: {template_path} (Index: {actual_template_index})")

    prs = Presentation(template_path)
    original_count = len(prs.slides)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    print(f"[PptxGen] Shablon slaydlar: {original_count}, Generatsiya: {len(slides)}")

    tmpl_slides = list(prs.slides)
    title_tmpl = tmpl_slides[0] if tmpl_slides else None
    content_tmpl = tmpl_slides[1] if len(tmpl_slides) > 1 else title_tmpl

    user_images = user_images or []

    for i, sd in enumerate(slides):
        if sd.slide_type == "title" and title_tmpl:
            new_slide = _clone_slide(prs, title_tmpl)
        elif content_tmpl:
            new_slide = _clone_slide(prs, content_tmpl)
        else:
            new_slide = prs.slides.add_slide(prs.slide_layouts[0])

        _inject_title(new_slide, sd.title)
        
        # Slayd rasmini aniqlaymiz
        img_path = None
        if user_images and i < len(user_images):
            img_path = user_images[i]
            
        # Dinamik Inch yordamida har bir elementni to'liq kafolatlangan chizish va ranglash
        _render_slide_content(new_slide, sd, img_path, slide_width, slide_height, actual_template_index)
        
        notes = getattr(sd, 'speaker_notes', '')  # speaker_notes maydoni bo'lmasa xato chiqmasin
        if notes:
            try:
                _inject_notes(new_slide, notes)
            except Exception:
                pass

    for i in range(original_count - 1, -1, -1):
        _delete_slide(prs, i)

    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    count = len(prs.slides)
    print(f"[PptxGen] Saqlandi: {output_path} ({count} slayd)")
    return output_path, count


async def generate_pptx(
    slides: list,
    output_path: str,
    style: str = "professional",
    template_index: Optional[int] = None,
    user_images: Optional[list] = None,
) -> tuple[str, int]:
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(
        None,
        _build_presentation,
        slides,
        output_path,
        template_index,
        user_images,
    )
    return result

PRO_TEMPLATES_DIR = Path(__file__).parent / "pro_templates"


def _get_pro_template_path(pro_design: Optional[str] = None, pro_design_variant: int = 1) -> str:
    """pro_design bo'limi va variant raqami bo'yicha PPTX faylini qaytaradi."""
    if pro_design:
        folder = PRO_TEMPLATES_DIR / pro_design
        if folder.exists():
            pptx_files = sorted(folder.glob("*.pptx"))
            if pptx_files:
                # variant 1 → birinchi fayl, variant 2 → ikkinchi fayl (agar mavjud bo'lsa)
                idx = min(pro_design_variant - 1, len(pptx_files) - 1)
                chosen = pptx_files[idx]
                print(f"[PptxGen] PRO shablon tanlandi: {chosen} (bo'lim={pro_design}, variant=#{pro_design_variant})")
                return str(chosen)

    # Fallback: barcha bo'limlar ichidan birorta .pptx topamiz
    for folder in PRO_TEMPLATES_DIR.iterdir():
        if folder.is_dir():
            pptx_files = sorted(folder.glob("*.pptx"))
            if pptx_files:
                print(f"[PptxGen] PRO shablon fallback: {pptx_files[0]}")
                return str(pptx_files[0])

    raise FileNotFoundError(f"PRO shablonlar topilmadi: {PRO_TEMPLATES_DIR}")


async def generate_pro_pptx(
    slides_data: dict,
    output_path: str,
    user_images: Optional[list] = None,
    pro_design: Optional[str] = None,
    pro_design_variant: int = 1,
    pro_plan_count: int = 5,
    pro_bibliography_type: str = "none",
) -> str:
    from pptx import Presentation
    import os
    import re
    TEMPLATE_PATH = _get_pro_template_path(pro_design, pro_design_variant)

    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(f"PRO Shablon topilmadi: {TEMPLATE_PATH}")
        
    print(f"[PptxGen] PRO Shablon ochilmoqda: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)
    
    # 1. ORTIQCHA SLAYDLAR VA ELEMENTLARNI TOZALASH
    slides_to_delete = []
    
    for idx, slide in enumerate(prs.slides):
        delete_this_slide = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_lower = shape.text.lower().replace("'", "_")
                
                # Agar slide reja_X_matni yoki xulosa matnini o'z ichiga olsa va X > pro_plan_count bo'lsa, slaydni o'chiramiz
                matches = re.findall(r'\{\{reja_(\d+)_matni', text_lower)
                for num_str in matches:
                    if int(num_str) > pro_plan_count:
                        delete_this_slide = True
                        break
                        
                matches_x = re.findall(r'\{\{reja_(\d+)_xulosa', text_lower)
                for num_str in matches_x:
                    if int(num_str) > pro_plan_count:
                        delete_this_slide = True
                        break

                if delete_this_slide:
                    break
                        
                # Adabiyotlar keraksiz bo'lsa, qidirib o'chiramiz
                if pro_bibliography_type == "none" and "{{adabiyotlar" in text_lower:
                    delete_this_slide = True
                    break
                    
        if delete_this_slide:
            slides_to_delete.append(idx)
            
    # Slaydlarni orqadan oldinga aylanib o'chiramiz (index siljib ketmasligi uchun)
    for idx in reversed(slides_to_delete):
        _delete_slide(prs, idx)

    # Qolgan slaydlarda (masalan asosiy Reja slaydi) oshiqcha punklarni o'chiramiz
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                for p_idx in reversed(range(len(paragraphs))):
                    p = paragraphs[p_idx]
                    p_text_lower = p.text.lower()
                    should_delete_p = False
                    
                    matches = re.findall(r'\{\{reja_(\d+)\}\}', p_text_lower)
                    for num_str in matches:
                        if int(num_str) > pro_plan_count:
                            should_delete_p = True
                            break
                            
                    if should_delete_p:
                        try:
                            p._element.getparent().remove(p._element)
                        except Exception:
                            p.text = ""


    
    image_replacements = {
        "ASOSIY RASM": 0,
        "PHOTO_2": 1,
        "PHOTO_3": 2,
        "PHOTO_4": 3,
        "PHOTO_5": 4
    }
    
    import re
    for slide in prs.slides:
        # Bir xil slaydda bitta teglarni necha marta takrorlanganini hisoblaymiz
        tag_counts = {}
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                shape_text = shape.text.replace("{{adabiyotlar_ro'yxati}}", "{{adabiyotlar_ro_yxati}}")
                for key in slides_data.keys():
                    tag = f"{{{{{key}}}}}"
                    occ = shape_text.count(tag)
                    if occ > 0:
                        tag_counts[key] = tag_counts.get(key, 0) + occ
        
        tag_current_index = {k: 0 for k in tag_counts}
        
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text and "}}" in run.text:
                            # Edge case ni to'g'irlash
                            run.text = run.text.replace("{{adabiyotlar_ro'yxati}}", "{{adabiyotlar_ro_yxati}}")
                            
                            for key, value in slides_data.items():
                                tag = f"{{{{{key}}}}}"
                                occ_in_run = run.text.count(tag)
                                for _ in range(occ_in_run):
                                    total_occ = tag_counts.get(key, 1)
                                    is_long = isinstance(value, str) and len(str(value)) > 60
                                    
                                    # Agar bitta slaydda bitta teg bir necha marta kelsa, matnni teng qismlarga bo'lamiz
                                    if total_occ > 1 and (is_long or str(key).endswith('_matni') or str(key).endswith('_xulosa')):
                                        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', str(value)) if s.strip()]
                                        idx = tag_current_index.get(key, 0)
                                        
                                        lst = sentences if len(sentences) >= total_occ else str(value).split()
                                        
                                        # Ro'yxatni total_occ ga mutanosib ravishda kesish
                                        k, m = divmod(len(lst), total_occ)
                                        start = idx * k + min(idx, m)
                                        end = (idx + 1) * k + min(idx + 1, m)
                                        chunk_items = lst[start:end]
                                        
                                        replacement_text = " ".join(chunk_items)
                                        tag_current_index[key] = idx + 1
                                    else:
                                        replacement_text = str(value)
                                        
                                    run.text = run.text.replace(tag, replacement_text, 1)
            
            pic_idx = None
            if hasattr(shape, "name"):
                pic_idx = image_replacements.get(shape.name.upper().strip())
            
            if pic_idx is None and shape.has_text_frame:
                content = shape.text.upper().strip()
                pic_idx = image_replacements.get(content)
            
            if pic_idx is not None and user_images and pic_idx < len(user_images) and user_images[pic_idx]:
                img_path = user_images[pic_idx]
                if os.path.exists(img_path):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    try:
                        sp = shape._element
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(img_path, left, top, width, height)
                    except Exception as e:
                        print(f"Rasm joylashda xato: {e}")

    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    final_slide_count = len(prs.slides)
    print(f"[PptxGen] PRO taqdimot saqlandi: {output_path} ({final_slide_count} slayd)")
    return output_path, final_slide_count

