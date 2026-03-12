"""
Taqdimot Generatsiya Pipeline (Professional Multi-Stage)
=========================================================
Oqim:
  1. Outline Generator  → slayd turlarini aniqlash
  2. Content Generator   → har bir slayd uchun mazmun yaratish (parallel)
  3. Image Fetcher       → Unsplash rasmlar (max 5)
  4. Template Engine     → HTML / PPTX / PDF yaratish
  5. Telegram            → fayllarni yuborish
"""
from __future__ import annotations

import asyncio
import json
import os
import uuid
from typing import Optional


# ──────────────────────────────────────────────────────────────
# 10 ta tayyor dizayn shablonlari (JSON fayllardan yuklanadi)
# ──────────────────────────────────────────────────────────────

def _load_design_templates() -> dict[int, dict]:
    import json
    from pathlib import Path

    templates = {}
    templates_dir = Path(__file__).parent / "templates"
    
    if templates_dir.exists():
        for file in templates_dir.glob("theme_*.json"):
            try:
                idx = int(file.stem.split("_")[1])
                with open(file, "r", encoding="utf-8") as f:
                    templates[idx] = json.load(f)
            except Exception as e:
                print(f"[Pipeline] Shablon yuklashda xato {file.name}: {e}")

    # Fallback to defaults if empty
    if not templates:
        templates[1] = {
            "page_bg": "#ffffff", "slide_bg": "#ffffff",
            "title_color": "#1a1a1a", "body_color": "#333333",
            "accent": "#c8e6c9", "num_color": "#999999",
            "font": "'Segoe UI', sans-serif", "theme_name": "Minimalist Macaron",
        }
    return templates

DESIGN_TEMPLATES: dict[int, dict] = _load_design_templates()

# Qo'llab-quvvatlanadigan slayd turlari
SLIDE_TYPES = [
    "title", "problem", "solution", "content", "chart",
    "image_text", "comparison", "quote", "team", "closing",
]

GROQ_MODEL = "llama-3.3-70b-versatile"


# ──────────────────────────────────────────────────────────────
# Pipeline
# ──────────────────────────────────────────────────────────────

class PresentationPipeline:
    """
    Professional multi-stage pipeline:
      1. _generate_outline()       → slayd turlari
      2. _generate_all_content()   → parallel content
      3. _fetch_images()           → Unsplash
      4. _build_pptx               → render
      5. _send_to_telegram         → yuborish
    """

    def _output_dir(self) -> str:
        d = os.path.join(os.path.expanduser("~"), "presentations_cache")
        os.makedirs(d, exist_ok=True)
        return d

    # ══════════════════════════════════════════════════════════
    # PUBLIC ENTRY POINT
    # ══════════════════════════════════════════════════════════

    async def run(
        self,
        topic: str,
        language: str,
        slide_count: int,
        style: str,
        extra_context: Optional[str],
        telegram_id: Optional[int],
        design_template: int = 1,
    ) -> dict:
        pid = str(uuid.uuid4())
        out = self._output_dir()

        # 1. Dizayn shabloni
        theme = DESIGN_TEMPLATES.get(design_template, DESIGN_TEMPLATES[1]).copy()
        print(f"[Pipeline] Shablon: {theme.get('theme_name', '?')}")

        # 2. Outline — slayd turlarini aniqlash
        outline = await self._generate_outline(topic, language, slide_count, extra_context)
        print(f"[Pipeline] Outline: {len(outline)} slayd turi aniqlandi")

        # 3. Har bir slayd uchun kontent yaratish (parallel)
        slides = await self._generate_all_content(outline, topic, language, extra_context)
        print(f"[Pipeline] Content: {len(slides)} slayd tayyor")

        # 4. Rasmlar (max 5, faqat image_prompt bor slaydlar)
        slides = await self._fetch_images(slides)

        # 5. PPTX
        pptx_path = f"{out}/{pid}.pptx"
        self._build_pptx(topic, slides, theme, pptx_path)

        # 6. Telegram
        tg_sent = False
        if telegram_id:
            tg_sent = await self._send_to_telegram(telegram_id, pptx_path, topic)

        return {
            "id": pid,
            "pptx_url": f"/api/v1/presentations/download/{pid}/pptx",
            "telegram_sent": tg_sent,
            "slide_count": len(slides),
        }

    # ══════════════════════════════════════════════════════════
    # STAGE 1: OUTLINE GENERATOR
    # ══════════════════════════════════════════════════════════

    async def _generate_outline(
        self, topic: str, language: str, slide_count: int, extra_context: Optional[str]
    ) -> list[dict]:
        """AI slayd turlarini aniqlaydi: title, problem, solution, ..."""
        from app.core.config import settings

        if not settings.GROQ_API_KEY:
            return self._default_outline(slide_count)

        extra = f"\nQo'shimcha kontekst: {extra_context}" if extra_context else ""
        types_str = ", ".join(SLIDE_TYPES)

        prompt = f"""Sen professional taqdimot strukturasini yaratuvchi AI yordamchisan.

Mavzu: {topic}{extra}
Slaydlar soni: {slide_count}

Quyidagi slayd turlaridan foydalanib, taqdimot strukturasini yarat:
{types_str}

QOIDALAR:
1. Birinchi slayd DOIM "title" turi bo'lsin
2. Oxirgi slayd DOIM "closing" turi bo'lsin
3. Mavzuga mos turlarni tanla
4. Har xil turlardan foydalanib, qiziqarli struktura yarat

Faqat JSON formatda javob ber:
[
  {{"type": "title", "hint": "Mavzu sarlavhasi va kirish"}},
  {{"type": "problem", "hint": "Muammo tavsifi"}},
  ...
]"""

        try:
            from groq import AsyncGroq
            client = AsyncGroq(api_key=settings.GROQ_API_KEY)

            resp = await client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[
                    {"role": "system", "content": "Faqat valid JSON array qaytarasan. Boshqa hech narsa yozma."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.6,
                max_tokens=1024,
            )

            raw = resp.choices[0].message.content.strip()
            raw = self._clean_json(raw)
            outline = json.loads(raw)

            # Validatsiya
            validated = []
            for item in outline:
                if isinstance(item, dict) and "type" in item:
                    stype = item["type"] if item["type"] in SLIDE_TYPES else "content"
                    validated.append({"type": stype, "hint": str(item.get("hint", ""))})

            return validated[:slide_count] if validated else self._default_outline(slide_count)

        except Exception as e:
            print(f"[Pipeline] Outline xatosi: {e}")
            return self._default_outline(slide_count)

    def _default_outline(self, count: int) -> list[dict]:
        """Fallback outline."""
        base = [
            {"type": "title", "hint": "Kirish"},
            {"type": "problem", "hint": "Muammo"},
            {"type": "solution", "hint": "Yechim"},
            {"type": "content", "hint": "Asosiy ma'lumot"},
            {"type": "chart", "hint": "Statistika"},
            {"type": "image_text", "hint": "Vizual"},
            {"type": "comparison", "hint": "Solishtirish"},
            {"type": "quote", "hint": "Iqtibos"},
            {"type": "team", "hint": "Jamoa"},
            {"type": "closing", "hint": "Xulosa"},
        ]
        result = base[:count]
        # Doim title bilan boshlash va closing bilan tugash
        if result[0]["type"] != "title":
            result[0] = {"type": "title", "hint": "Kirish"}
        if len(result) > 1 and result[-1]["type"] != "closing":
            result[-1] = {"type": "closing", "hint": "Xulosa"}
        return result

    # ══════════════════════════════════════════════════════════
    # STAGE 2: CONTENT GENERATOR (PARALLEL)
    # ══════════════════════════════════════════════════════════

    async def _generate_all_content(
        self, outline: list[dict], topic: str, language: str, extra_context: Optional[str]
    ) -> list[dict]:
        """Barcha slaydlar uchun parallel content yaratish."""
        tasks = [
            self._generate_slide_content(item, topic, language, extra_context, i, len(outline))
            for i, item in enumerate(outline)
        ]
        slides = await asyncio.gather(*tasks)
        return list(slides)

    async def _generate_slide_content(
        self,
        outline_item: dict,
        topic: str,
        language: str,
        extra_context: Optional[str],
        index: int,
        total: int,
    ) -> dict:
        """Bitta slayd uchun kontent yaratish."""
        from app.core.config import settings

        slide_type = outline_item["type"]
        hint = outline_item.get("hint", "")

        if not settings.GROQ_API_KEY:
            return self._demo_slide(slide_type, topic, hint, index)

        lang_map = {"uz": "O'zbek tilida", "ru": "на русском языке", "en": "in English"}
        lang_str = lang_map.get(language, "in English")
        extra = f"\nQo'shimcha kontekst: {extra_context}" if extra_context else ""

        # Slayd turiga qarab maxsus instruktsiya
        type_instructions = {
            "title": "Kirish slayd. Katta sarlavha va qisqa tavsif (1-2 jumla subtitle). Body kerak emas.",
            "problem": "Muammo slayd. 50-80 so'zlik akademik matn muammoni tavsiflaydi. 3-4 ta asosiy muammo nuqtalari.",
            "solution": "Yechim slayd. 50-80 so'zlik akademik matn yechimni tavsifla. 3-4 ta yechim yo'llari.",
            "content": "Asosiy kontent slayd. 60-100 so'zlik chuqur akademik matn. Ma'lumotli paragraf.",
            "chart": "Statistika slayd. Raqamlar va dalillarga boy matn. Kamida 3-4 ta raqamli ma'lumot (points da).",
            "image_text": "Vizual tushuntirish slayd. 50-80 so'zlik matn + image_prompt berish SHART.",
            "comparison": "Solishtirish slayd. 2 ta narsani solishtiradigan matn. Points da har birining afzalliklari.",
            "quote": "Iqtibos slayd. Mashhur shaxs yoki manba dan iqtibos. body da iqtibos, subtitle da muallif.",
            "team": "Jamoa/tashkilot slayd. Ishtirokchilar yoki tashkilotlar haqida. Points da ro'yxat.",
            "closing": "Xulosa slayd. Asosiy xulosalar va keyingi qadamlar. 3-4 yakuniy fikr.",
        }

        instruction = type_instructions.get(slide_type, "Akademik matn yozing. 50-100 so'z.")

        # image_prompt faqat image_text, problem, solution, content turlarda berish mumkin
        image_note = ""
        if slide_type in ("image_text", "problem", "solution", "content", "title"):
            image_note = '\n"image_prompt" maydonini ham qo\'shing — Unsplash dan rasm izlash uchun ingliz tilida 2-4 so\'zlik tavsif.'

        prompt = f"""Taqdimot slayd kontenti yarat.
Mavzu: {topic}
Slayd turi: {slide_type}
Hint: {hint}
Bu {index + 1}-slayd (jami {total} ta)
Til: {lang_str}{extra}

{instruction}

JSON formatda javob ber:
{{
  "title": "Slayd sarlavhasi",
  "subtitle": "Qo'shimcha sarlavha (ixtiyoriy)",
  "body": "Akademik paragraf matni (50-100 so'z)",
  "points": ["Nuqta 1", "Nuqta 2", "Nuqta 3"],
  "emoji": "🎯"{image_note}
}}

MUHIM: body maydoni 50-100 so'zdan iborat bo'lsin. Points maydoni 3-5 ta qisqa gap."""

        try:
            from groq import AsyncGroq
            client = AsyncGroq(api_key=settings.GROQ_API_KEY)

            resp = await client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[
                    {"role": "system", "content": "Faqat valid JSON object qaytarasan. Boshqa hech narsa yozma."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.7,
                max_tokens=1024,
            )

            raw = resp.choices[0].message.content.strip()
            raw = self._clean_json(raw)

            # { ... } qismini olish
            start = raw.find("{")
            end = raw.rfind("}") + 1
            if start != -1 and end > 0:
                raw = raw[start:end]

            try:
                data = json.loads(raw)
            except Exception as parse_e:
                with open("json_err.log", "w", encoding="utf-8") as fe:
                    fe.write(f"Error: {parse_e}\nRaw:\n{raw}")
                raise parse_e

            slide = {
                "type": slide_type,
                "title": str(data.get("title", hint or topic)),
                "subtitle": str(data.get("subtitle", "")),
                "body": str(data.get("body", "")),
                "points": [str(p) for p in data.get("points", [])],
                "emoji": str(data.get("emoji", "📌")),
            }
            if data.get("image_prompt"):
                slide["image_prompt"] = str(data["image_prompt"])

            return slide

        except Exception as e:
            print(f"[Pipeline] Content xatosi (slayd {index + 1}, {slide_type}): {e}")
            return self._demo_slide(slide_type, topic, hint, index)

    def _demo_slide(self, slide_type: str, topic: str, hint: str, index: int) -> dict:
        """Fallback demo slide."""
        emojis = {"title": "🚀", "problem": "⚠️", "solution": "💡", "content": "📖",
                  "chart": "📊", "image_text": "🖼️", "comparison": "⚖️",
                  "quote": "💬", "team": "👥", "closing": "🎯"}
        return {
            "type": slide_type,
            "title": f"{topic}" if slide_type == "title" else hint or slide_type.capitalize(),
            "subtitle": f"{topic} haqida" if slide_type == "title" else "",
            "body": f"{topic} haqida {index + 1}-qism ma'lumoti. Bu demo matn.",
            "points": ["Muhim nuqta 1", "Muhim nuqta 2", "Muhim nuqta 3"],
            "emoji": emojis.get(slide_type, "📌"),
        }

    # ══════════════════════════════════════════════════════════
    # STAGE 3: IMAGE FETCHER
    # ══════════════════════════════════════════════════════════

    async def _fetch_images(self, slides: list[dict]) -> list[dict]:
        """Taqdimot uchun rasmlar olish (max 5)."""
        image_count = 0
        for slide in slides:
            if slide.get("image_prompt") and image_count < 5:
                # Use a reliable free image generation API instead of deprecated Unsplash Source
                import urllib.parse
                q = urllib.parse.quote(slide["image_prompt"])
                slide["image_url"] = f"https://image.pollinations.ai/prompt/{q}?width=800&height=400&nologo=true"
                image_count += 1
            elif "image_prompt" in slide and image_count >= 5:
                del slide["image_prompt"]  # Limit exceeded
        print(f"[Pipeline] Rasmlar: {image_count} ta rasm tayyor")
        return slides
    # ══════════════════════════════════════════════════════════
    # STAGE 4: PPTX BUILDER
    # ══════════════════════════════════════════════════════════

    def _build_pptx(self, topic: str, slides: list[dict], theme: dict, path: str) -> None:
        try:
            import httpx
            import tempfile
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            cfg = theme

            def hex_to_rgb(h: str) -> RGBColor:
                h = h.lstrip("#")
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            slide_bg_rgb = hex_to_rgb(cfg["slide_bg"])
            title_rgb = hex_to_rgb(cfg["title_color"])
            body_rgb = hex_to_rgb(cfg["body_color"])
            accent_rgb = hex_to_rgb(cfg["accent"])

            for slide_data in slides:
                layout = prs.slide_layouts[6]  # blank
                sl = prs.slides.add_slide(layout)

                # Fon
                bg = sl.background.fill
                bg.solid()
                bg.fore_color.rgb = slide_bg_rgb

                # Accent chiziq (chapdan)
                bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
                bar.fill.solid()
                bar.fill.fore_color.rgb = accent_rgb
                bar.line.fill.background()

                stype = slide_data.get("type", "content")

                # Emoji
                tx = sl.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(2), Inches(0.8))
                tf = tx.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.get("emoji", "📌")
                p.font.size = Pt(32)

                # Title
                title_top = Inches(1.4)
                if stype == "title":
                    # Title slayd — markazda, katta
                    tx = sl.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2))
                    tf = tx.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_data.get("title", "")
                    p.font.bold = True
                    p.font.size = Pt(44)
                    p.font.color.rgb = title_rgb
                    p.alignment = PP_ALIGN.CENTER

                    # Subtitle
                    subtitle = slide_data.get("subtitle", "")
                    if subtitle:
                        p2 = tf.add_paragraph()
                        p2.text = subtitle
                        p2.font.size = Pt(24)
                        p2.font.color.rgb = accent_rgb
                        p2.alignment = PP_ALIGN.CENTER
                        p2.space_before = Pt(16)
                else:
                    # Normal title
                    tx = sl.shapes.add_textbox(Inches(0.5), title_top, Inches(12.5), Inches(1.2))
                    tf = tx.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_data.get("title", "")
                    p.font.bold = True
                    p.font.size = Pt(32)
                    p.font.color.rgb = title_rgb

                    # Subtitle
                    subtitle = slide_data.get("subtitle", "")
                    if subtitle:
                        p2 = tf.add_paragraph()
                        p2.text = subtitle
                        p2.font.size = Pt(20)
                        p2.font.color.rgb = accent_rgb
                        p2.space_before = Pt(6)

                # Skip body/points for title slide
                if stype == "title":
                    continue

                # Rasm uchun joy aniqlash
                has_image = bool(slide_data.get("image_url"))
                body_left = Inches(0.5)
                body_width = Inches(12.5)
                body_top = Inches(2.8)

                if has_image:
                    try:
                        img_url = slide_data["image_url"]
                        with httpx.Client(timeout=15.0, follow_redirects=True) as http:
                            resp = http.get(img_url)
                        if resp.status_code == 200 and len(resp.content) > 1000:
                            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                                tmp.write(resp.content)
                                tmp_path = tmp.name
                            if stype == "image_text":
                                sl.shapes.add_picture(tmp_path, Inches(0.5), Inches(2.8), Inches(5.5), Inches(4.0))
                                body_left = Inches(6.5)
                                body_width = Inches(6.5)
                            else:
                                sl.shapes.add_picture(tmp_path, Inches(7.5), Inches(1.5), Inches(5.5), Inches(3.5))
                                body_width = Inches(6.8)
                            os.remove(tmp_path)
                    except Exception as img_err:
                        print(f"[Pipeline] PPTX rasm xatosi: {img_err}")

                # Body matn
                body = slide_data.get("body", "")
                if body:
                    tx = sl.shapes.add_textbox(body_left, body_top, body_width, Inches(2.0))
                    tf = tx.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = body
                    p.font.size = Pt(16)
                    p.font.color.rgb = body_rgb
                    p.line_spacing = Pt(24)

                # Points
                points = slide_data.get("points", [])
                if points:
                    points_top = body_top + Inches(2.2) if body else body_top
                    tx = sl.shapes.add_textbox(body_left, points_top, body_width, Inches(3.0))
                    tf = tx.text_frame
                    tf.word_wrap = True
                    first = True
                    for point in points:
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        p.text = f"• {point}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = body_rgb
                        p.space_before = Pt(6)

            prs.save(path)

        except ImportError:
            print("[Pipeline] python-pptx o'rnatilmagan!")
            open(path, "wb").close()
        except Exception as e:
            print(f"[Pipeline] PPTX xatosi: {e}")
            import traceback
            traceback.print_exc()
            open(path, "wb").close()
    # ══════════════════════════════════════════════════════════
    # STAGE 5: TELEGRAM
    # ══════════════════════════════════════════════════════════

    async def _send_to_telegram(self, telegram_id: int, pptx_path: str, topic: str) -> bool:
        try:
            import httpx
            from app.core.config import settings

            api_url = f"https://api.telegram.org/bot{settings.TELEGRAM_BOT_TOKEN}"
            sent_any = False

            async with httpx.AsyncClient(timeout=60.0) as client:
                # PPTX
                if os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 0:
                    with open(pptx_path, "rb") as fh:
                        resp = await client.post(
                            f"{api_url}/sendDocument",
                            data={
                                "chat_id": str(telegram_id),
                                "caption": f"📊 <b>{topic}</b>\n\n✅ PPTX taqdimotingiz tayyor!",
                                "parse_mode": "HTML",
                            },
                            files={"document": (f"{topic[:40]}.pptx", fh, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                        )
                    if resp.status_code == 200:
                        sent_any = True

            return sent_any

        except Exception as e:
            print(f"[Pipeline] Telegram yuborish xatosi: {e}")
            return False

    # ══════════════════════════════════════════════════════════
    # HELPERS
    # ══════════════════════════════════════════════════════════

    @staticmethod
    def _esc(text: str) -> str:
        return (text
                .replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;"))

    @staticmethod
    def _clean_json(raw: str) -> str:
        if "```json" in raw:
            raw = raw.split("```json")[1].split("```")[0].strip()
        elif "```" in raw:
            raw = raw.split("```")[1].split("```")[0].strip()
            
        start_obj = raw.find("{")
        start_arr = raw.find("[")
        if start_obj == -1:
            start = start_arr
        elif start_arr == -1:
            start = start_obj
        else:
            start = min(start_obj, start_arr)
            
        end_obj = raw.rfind("}")
        end_arr = raw.rfind("]")
        end = max(end_obj, end_arr) + 1
        
        if start != -1 and end > 0:
            raw = raw[start:end]
        return raw.strip()
