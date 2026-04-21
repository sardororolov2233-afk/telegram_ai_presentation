import sys
import collections
from pptx import Presentation

def extract_tags(pptx_path):
    prs = Presentation(pptx_path)
    tags = collections.defaultdict(list)
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                if '{' in text and '}' in text:
                    # simplistic extraction
                    words = text.split()
                    for w in words:
                        if '{' in w and '}' in w:
                            tags[slide_num].append(w)
    
    for slide_num, tgs in tags.items():
        print(f"Slide {slide_num}: {', '.join(tgs)}")

if __name__ == "__main__":
    extract_tags(r"d:\Новая папка\backend\app\services\presentation\pro_templates\fizika\Albert Einstein Slides, копия.pptx")
